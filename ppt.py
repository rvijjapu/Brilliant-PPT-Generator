"""
Brilliant PPT Generator - Ultimate Edition
AI-Powered Professional Presentations | Google Gemini FREE
"""

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import io
import json
import re
import requests
import time

# ============ CONFIGURATION ============
GEMINI_API_KEY = "AIzaSyC7IYzC4zdlKVRdA5YAInwN7vsPnypIin4"

st.set_page_config(
    page_title="Brilliant PPT Generator",
    page_icon="✨",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# ============ MODERN PROFESSIONAL THEMES ============
THEMES = {
    'Corporate Blue': {
        'header_bg': (37, 99, 235),
        'accent': (59, 130, 246),
        'secondary': (239, 246, 255),
        'body_text': (30, 41, 59),
        'gradient': 'linear-gradient(135deg, #2563EB 0%, #3B82F6 100%)',
        'preview_dark': '#2563EB'
    },
    'Emerald Pro': {
        'header_bg': (16, 185, 129),
        'accent': (52, 211, 153),
        'secondary': (236, 253, 245),
        'body_text': (30, 41, 59),
        'gradient': 'linear-gradient(135deg, #10B981 0%, #34D399 100%)',
        'preview_dark': '#10B981'
    },
    'Sunset Orange': {
        'header_bg': (249, 115, 22),
        'accent': (251, 146, 60),
        'secondary': (255, 247, 237),
        'body_text': (30, 41, 59),
        'gradient': 'linear-gradient(135deg, #F97316 0%, #FB923C 100%)',
        'preview_dark': '#F97316'
    },
    'Royal Violet': {
        'header_bg': (139, 92, 246),
        'accent': (167, 139, 250),
        'secondary': (245, 243, 255),
        'body_text': (30, 41, 59),
        'gradient': 'linear-gradient(135deg, #8B5CF6 0%, #A78BFA 100%)',
        'preview_dark': '#8B5CF6'
    },
    'Rose Pink': {
        'header_bg': (244, 63, 94),
        'accent': (251, 113, 133),
        'secondary': (255, 241, 242),
        'body_text': (30, 41, 59),
        'gradient': 'linear-gradient(135deg, #F43F5E 0%, #FB7185 100%)',
        'preview_dark': '#F43F5E'
    },
    'Ocean Teal': {
        'header_bg': (20, 184, 166),
        'accent': (45, 212, 191),
        'secondary': (240, 253, 250),
        'body_text': (30, 41, 59),
        'gradient': 'linear-gradient(135deg, #14B8A6 0%, #2DD4BF 100%)',
        'preview_dark': '#14B8A6'
    },
    'Golden Amber': {
        'header_bg': (245, 158, 11),
        'accent': (251, 191, 36),
        'secondary': (255, 251, 235),
        'body_text': (30, 41, 59),
        'gradient': 'linear-gradient(135deg, #F59E0B 0%, #FBBF24 100%)',
        'preview_dark': '#F59E0B'
    },
    'Sky Blue': {
        'header_bg': (14, 165, 233),
        'accent': (56, 189, 248),
        'secondary': (240, 249, 255),
        'body_text': (30, 41, 59),
        'gradient': 'linear-gradient(135deg, #0EA5E9 0%, #38BDF8 100%)',
        'preview_dark': '#0EA5E9'
    },
    'Slate Pro': {
        'header_bg': (71, 85, 105),
        'accent': (100, 116, 139),
        'secondary': (248, 250, 252),
        'body_text': (30, 41, 59),
        'gradient': 'linear-gradient(135deg, #475569 0%, #64748B 100%)',
        'preview_dark': '#475569'
    },
    'Indigo Night': {
        'header_bg': (99, 102, 241),
        'accent': (129, 140, 248),
        'secondary': (238, 242, 255),
        'body_text': (30, 41, 59),
        'gradient': 'linear-gradient(135deg, #6366F1 0%, #818CF8 100%)',
        'preview_dark': '#6366F1'
    }
}


# ============ AI CONTENT GENERATION ============
def generate_presentation_with_ai(title: str, content: str, num_slides: int = 6) -> dict:
    """Generate professional presentation using Google Gemini AI"""
    
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={GEMINI_API_KEY}"
    
    prompt = f"""You are a world-class presentation designer who creates CEO-level, award-winning PowerPoint presentations.

TASK: Create a compelling, professional presentation from the user's input. Transform raw ideas into a powerful narrative.

PRESENTATION TITLE: {title}

USER'S RAW CONTENT:
{content}

CRITICAL INSTRUCTIONS:
1. CREATE A STORY: Build a clear narrative arc - Opening Hook → Context → Main Points → Evidence → Conclusion → Call to Action
2. PROFESSIONAL LANGUAGE: Use executive-level, impactful business language
3. FIX EVERYTHING: Correct all spelling, grammar, and unclear phrases
4. ENHANCE WEAK CONTENT: If input is vague, incomplete, or even nonsensical - intelligently create relevant professional content based on the title
5. HANDLE ANY INPUT: Even if content seems random or wrong, create a meaningful presentation around the title theme
6. STRUCTURE: Generate exactly {num_slides} content slides (excluding title and thank you slides)

SLIDE STRUCTURE:
- Each slide: Compelling title (4-8 words) + 3-5 bullet points
- Bullets: Clear, concise, action-oriented (max 15 words each)
- Flow: Each slide should logically lead to the next

SLIDE TYPES TO INCLUDE:
1. Executive Summary / Opening Hook
2. Problem / Opportunity / Context
3. Solution / Approach / Key Points
4. Benefits / Value Proposition / Results
5. Implementation / Timeline / Process
6. Conclusion / Next Steps / Call to Action

OUTPUT FORMAT - Return ONLY valid JSON, no markdown, no explanation:
{{
    "title": "{title}",
    "slides": [
        {{
            "title": "Engaging Slide Title",
            "bullets": [
                "First impactful point with clear value",
                "Second point building the narrative",
                "Third point with supporting evidence",
                "Fourth point driving action"
            ]
        }}
    ]
}}

IMPORTANT: You MUST generate {num_slides} slides with meaningful content. Never return empty slides.
Generate the complete JSON now:"""

    headers = {"Content-Type": "application/json"}
    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {
            "temperature": 0.8,
            "topK": 40,
            "topP": 0.95,
            "maxOutputTokens": 4096
        }
    }
    
    try:
        response = requests.post(url, headers=headers, json=payload, timeout=60)
        response.raise_for_status()
        
        result = response.json()
        text = result['candidates'][0]['content']['parts'][0]['text']
        
        # Clean response
        text = text.strip()
        if '```json' in text:
            text = text.split('```json')[1]
        if '```' in text:
            text = text.split('```')[0]
        text = text.strip()
        
        # Parse JSON
        data = json.loads(text)
        
        # Validate slides
        if 'slides' not in data or len(data['slides']) == 0:
            return None
        
        # Ensure all slides have content
        valid_slides = []
        for slide in data['slides']:
            if slide.get('title') and slide.get('bullets') and len(slide['bullets']) > 0:
                # Clean bullets
                clean_bullets = [b.strip() for b in slide['bullets'] if b and b.strip()]
                if clean_bullets:
                    valid_slides.append({
                        'title': slide['title'].strip(),
                        'bullets': clean_bullets[:5]  # Max 5 bullets
                    })
        
        if valid_slides:
            data['slides'] = valid_slides
            return data
        
        return None
        
    except json.JSONDecodeError as e:
        st.warning(f"JSON parsing issue, retrying...")
        return None
    except requests.exceptions.RequestException as e:
        st.error(f"API Error: {str(e)}")
        return None
    except Exception as e:
        st.error(f"Error: {str(e)}")
        return None


def smart_fallback_parser(title: str, content: str) -> dict:
    """Intelligent fallback parser when AI fails"""
    
    lines = [l.strip() for l in content.split('\n') if l.strip()]
    
    # If no content, generate based on title
    if not lines or len(content.strip()) < 10:
        return generate_template_presentation(title)
    
    slides = []
    current_slide = None
    
    for line in lines:
        # Detect bullet points
        is_bullet = bool(re.match(r'^[-*•→>✓✔◆■□▪▸►]\s*', line)) or \
                    bool(re.match(r'^\d+[.)]\s*', line))
        
        if is_bullet:
            bullet_text = re.sub(r'^[-*•→>✓✔◆■□▪▸►\d.)]+\s*', '', line).strip()
            if bullet_text:
                if not current_slide:
                    current_slide = {'title': 'Key Highlights', 'bullets': []}
                if len(current_slide['bullets']) < 5:
                    current_slide['bullets'].append(bullet_text)
        else:
            # It's a heading/title
            if current_slide and current_slide.get('bullets'):
                slides.append(current_slide)
            current_slide = {'title': line[:60], 'bullets': []}
    
    # Add last slide
    if current_slide:
        if current_slide.get('bullets'):
            slides.append(current_slide)
        elif current_slide.get('title'):
            current_slide['bullets'] = [
                f"Overview of {current_slide['title']}",
                "Key points and considerations",
                "Strategic importance"
            ]
            slides.append(current_slide)
    
    # If still no slides, create from content chunks
    if not slides:
        slides = create_slides_from_text(title, lines)
    
    # Ensure minimum slides
    if len(slides) < 3:
        slides = ensure_minimum_slides(title, slides)
    
    return {"title": title, "slides": slides}


def generate_template_presentation(title: str) -> dict:
    """Generate a template presentation based on title"""
    
    return {
        "title": title,
        "slides": [
            {
                "title": "Executive Summary",
                "bullets": [
                    f"Comprehensive overview of {title}",
                    "Key objectives and strategic goals",
                    "Expected outcomes and deliverables",
                    "Timeline and resource requirements"
                ]
            },
            {
                "title": "Current Landscape",
                "bullets": [
                    "Analysis of the present situation",
                    "Challenges and opportunities identified",
                    "Market trends and competitive factors",
                    "Stakeholder perspectives and needs"
                ]
            },
            {
                "title": "Strategic Approach",
                "bullets": [
                    "Recommended methodology and framework",
                    "Key initiatives and action items",
                    "Resource allocation strategy",
                    "Risk mitigation measures"
                ]
            },
            {
                "title": "Value Proposition",
                "bullets": [
                    "Tangible benefits and ROI potential",
                    "Competitive advantages gained",
                    "Efficiency improvements expected",
                    "Long-term strategic value"
                ]
            },
            {
                "title": "Implementation Roadmap",
                "bullets": [
                    "Phase 1: Foundation and planning",
                    "Phase 2: Execution and deployment",
                    "Phase 3: Optimization and scaling",
                    "Key milestones and checkpoints"
                ]
            },
            {
                "title": "Next Steps & Recommendations",
                "bullets": [
                    "Immediate action items to initiate",
                    "Decision points requiring approval",
                    "Resource and budget requirements",
                    "Success metrics and KPIs"
                ]
            }
        ]
    }


def create_slides_from_text(title: str, lines: list) -> list:
    """Create slides from raw text lines"""
    
    slides = []
    chunk_size = 4
    
    for i in range(0, len(lines), chunk_size):
        chunk = lines[i:i + chunk_size]
        
        # First non-bullet line is title
        slide_title = None
        bullets = []
        
        for line in chunk:
            clean_line = re.sub(r'^[-*•→>✓✔◆■□▪▸►\d.)]+\s*', '', line).strip()
            if not slide_title and not re.match(r'^[-*•]', line):
                slide_title = clean_line[:60]
            elif clean_line:
                bullets.append(clean_line)
        
        if not slide_title:
            slide_title = f"Key Points {len(slides) + 1}"
        
        if not bullets:
            bullets = [f"Details about {slide_title}"]
        
        slides.append({
            'title': slide_title,
            'bullets': bullets[:5]
        })
    
    return slides


def ensure_minimum_slides(title: str, existing_slides: list) -> list:
    """Ensure presentation has minimum required slides"""
    
    default_slides = [
        {
            "title": "Overview",
            "bullets": [
                f"Introduction to {title}",
                "Key objectives and goals",
                "Scope and deliverables"
            ]
        },
        {
            "title": "Key Insights",
            "bullets": [
                "Primary findings and analysis",
                "Critical success factors",
                "Strategic considerations"
            ]
        },
        {
            "title": "Recommendations",
            "bullets": [
                "Suggested next steps",
                "Implementation priorities",
                "Expected outcomes"
            ]
        }
    ]
    
    while len(existing_slides) < 3:
        if default_slides:
            existing_slides.append(default_slides.pop(0))
        else:
            break
    
    return existing_slides


# ============ POWERPOINT GENERATION ============
def create_pptx(title: str, slides_data: list, theme_name: str) -> io.BytesIO:
    """Create professional PowerPoint presentation"""
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)
    
    theme = THEMES.get(theme_name, THEMES['Corporate Blue'])
    
    # ===== TITLE SLIDE =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Gradient-like background (solid color)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(*theme['header_bg'])
    
    # Decorative accent bar
    accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor(*theme['accent'])
    accent_bar.line.fill.background()
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.0), Inches(11.833), Inches(2.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(50)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Calibri Light'
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.5), Inches(11.833), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Professional Presentation"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(*theme['secondary'])
    p.font.name = 'Calibri'
    
    # ===== CONTENT SLIDES =====
    for idx, slide_content in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        
        # White background
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(255, 255, 255)
        
        # Header bar
        header = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.25))
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(*theme['header_bg'])
        header.line.fill.background()
        
        # Slide title
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.28), Inches(11.833), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_content.get('title', f'Slide {idx + 1}')[:60]
        p.font.size = Pt(34)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.font.name = 'Calibri Light'
        
        # Vertical accent bar
        accent = slide.shapes.add_shape(1, Inches(0.75), Inches(1.55), Inches(0.12), Inches(5.0))
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(*theme['accent'])
        accent.line.fill.background()
        
        # Bullet points
        bullets = slide_content.get('bullets', [])
        if bullets:
            content_box = slide.shapes.add_textbox(Inches(1.25), Inches(1.7), Inches(11.333), Inches(5.3))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for i, bullet in enumerate(bullets[:5]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = f"•   {bullet}"
                p.font.size = Pt(21)
                p.font.color.rgb = RGBColor(*theme['body_text'])
                p.font.name = 'Calibri'
                p.space_before = Pt(8)
                p.space_after = Pt(12)
        
        # Slide number
        num_box = slide.shapes.add_textbox(Inches(12.4), Inches(7.0), Inches(0.7), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(idx + 1)
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(160, 160, 160)
        p.font.name = 'Calibri'
    
    # ===== THANK YOU SLIDE =====
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(*theme['header_bg'])
    
    # Decorative bar
    accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(13.333), Inches(0.08))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = RGBColor(*theme['accent'])
    accent_bar.line.fill.background()
    
    # Thank you text
    thanks_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.2), Inches(11.833), Inches(1.5))
    tf = thanks_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Thank You"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Calibri Light'
    
    # Questions text
    q_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.2), Inches(11.833), Inches(0.8))
    tf = q_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Questions & Discussion"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(26)
    p.font.color.rgb = RGBColor(*theme['secondary'])
    p.font.name = 'Calibri'
    
    # Save
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


# ============ MODERN UI STYLES ============
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Plus+Jakarta+Sans:wght@400;500;600;700;800&display=swap');

* {
    font-family: 'Plus Jakarta Sans', sans-serif;
}

.stApp {
    background: linear-gradient(135deg, #F8FAFC 0%, #EEF2FF 50%, #F0FDF4 100%);
}

/* Hide default elements */
#MainMenu, footer, header {visibility: hidden;}
.stDeployButton {display: none;}

/* Main container */
.block-container {
    padding: 1rem 2rem 2rem 2rem;
    max-width: 1300px;
}

/* Hero Header */
.hero-section {
    text-align: center;
    padding: 1.5rem 0 2rem 0;
    margin-bottom: 1rem;
}

.hero-title {
    font-size: 2.75rem;
    font-weight: 800;
    background: linear-gradient(135deg, #2563EB 0%, #7C3AED 50%, #EC4899 100%);
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
    margin: 0;
    letter-spacing: -0.02em;
}

.hero-subtitle {
    font-size: 1.1rem;
    color: #64748B;
    margin-top: 0.5rem;
    font-weight: 500;
}

.hero-badge {
    display: inline-flex;
    align-items: center;
    gap: 0.4rem;
    background: linear-gradient(135deg, #10B981, #059669);
    color: white;
    padding: 0.35rem 1rem;
    border-radius: 50px;
    font-size: 0.8rem;
    font-weight: 600;
    margin-top: 0.75rem;
    box-shadow: 0 4px 15px rgba(16, 185, 129, 0.3);
}

/* Cards */
.modern-card {
    background: white;
    border-radius: 20px;
    padding: 1.5rem;
    box-shadow: 0 4px 25px rgba(0, 0, 0, 0.06);
    border: 1px solid rgba(226, 232, 240, 0.8);
    margin-bottom: 1rem;
    transition: transform 0.2s, box-shadow 0.2s;
}

.modern-card:hover {
    box-shadow: 0 8px 35px rgba(0, 0, 0, 0.1);
}

.card-header {
    display: flex;
    align-items: center;
    gap: 0.6rem;
    margin-bottom: 1rem;
    padding-bottom: 0.75rem;
    border-bottom: 1px solid #F1F5F9;
}

.card-header-icon {
    font-size: 1.25rem;
}

.card-header-text {
    font-size: 1rem;
    font-weight: 700;
    color: #1E293B;
}

/* Theme Grid */
.theme-grid {
    display: grid;
    grid-template-columns: repeat(5, 1fr);
    gap: 0.6rem;
}

.theme-item {
    text-align: center;
    cursor: pointer;
    padding: 0.4rem;
    border-radius: 10px;
    transition: transform 0.15s;
}

.theme-item:hover {
    transform: scale(1.05);
}

.theme-preview-box {
    height: 32px;
    border-radius: 8px;
    margin-bottom: 0.3rem;
}

.theme-label {
    font-size: 0.65rem;
    color: #64748B;
    font-weight: 500;
}

/* Slide Preview */
.slide-card {
    background: white;
    border-radius: 14px;
    overflow: hidden;
    box-shadow: 0 6px 25px rgba(0, 0, 0, 0.12);
    margin: 0.6rem 0;
    border: 1px solid #E2E8F0;
}

.slide-card-header {
    padding: 0.85rem 1.25rem;
    color: white;
    font-weight: 600;
    font-size: 1rem;
}

.slide-card-body {
    padding: 1rem 1.25rem;
    background: white;
    min-height: 100px;
}

.slide-card-body p {
    margin: 0.35rem 0;
    color: #334155;
    font-size: 0.9rem;
    line-height: 1.5;
}

/* Progress indicator */
.progress-step {
    display: flex;
    align-items: center;
    gap: 0.5rem;
    padding: 0.6rem 1rem;
    background: #F8FAFC;
    border-radius: 10px;
    margin: 0.4rem 0;
    border-left: 3px solid #10B981;
}

.progress-step-text {
    font-size: 0.9rem;
    color: #475569;
}

/* Stats badges */
.stats-row {
    display: flex;
    gap: 0.75rem;
    margin: 0.75rem 0;
    flex-wrap: wrap;
}

.stat-badge {
    background: linear-gradient(135deg, #F8FAFC, #EEF2FF);
    padding: 0.5rem 1rem;
    border-radius: 10px;
    font-size: 0.85rem;
    color: #475569;
    font-weight: 500;
    border: 1px solid #E2E8F0;
}

/* Feature list */
.feature-list {
    background: linear-gradient(135deg, #F0FDF4 0%, #ECFDF5 100%);
    border-radius: 12px;
    padding: 1rem;
    margin: 0.75rem 0;
    border: 1px solid #BBF7D0;
}

.feature-item {
    display: flex;
    align-items: center;
    gap: 0.5rem;
    padding: 0.3rem 0;
    font-size: 0.85rem;
    color: #166534;
}

/* Buttons */
.stButton > button {
    width: 100%;
    padding: 0.9rem 1.5rem;
    font-size: 1rem;
    font-weight: 700;
    border-radius: 14px;
    border: none;
    transition: all 0.2s;
    letter-spacing: 0.02em;
}

.stButton > button:hover {
    transform: translateY(-3px);
    box-shadow: 0 8px 25px rgba(0, 0, 0, 0.15);
}

.stDownloadButton > button {
    width: 100%;
    padding: 1rem 1.5rem;
    font-size: 1rem;
    font-weight: 700;
    border-radius: 14px;
    background: linear-gradient(135deg, #10B981 0%, #059669 100%);
    color: white;
    border: none;
    box-shadow: 0 4px 15px rgba(16, 185, 129, 0.3);
}

.stDownloadButton > button:hover {
    transform: translateY(-2px);
    box-shadow: 0 8px 25px rgba(16, 185, 129, 0.4);
}

/* Input styling */
.stTextInput > div > div > input {
    border-radius: 12px;
    border: 2px solid #E2E8F0;
    padding: 0.75rem 1rem;
    font-size: 1rem;
    transition: border-color 0.2s, box-shadow 0.2s;
}

.stTextInput > div > div > input:focus {
    border-color: #2563EB;
    box-shadow: 0 0 0 3px rgba(37, 99, 235, 0.1);
}

.stTextArea > div > div > textarea {
    border-radius: 12px;
    border: 2px solid #E2E8F0;
    padding: 1rem;
    font-size: 0.95rem;
    line-height: 1.6;
    transition: border-color 0.2s, box-shadow 0.2s;
}

.stTextArea > div > div > textarea:focus {
    border-color: #2563EB;
    box-shadow: 0 0 0 3px rgba(37, 99, 235, 0.1);
}

/* Success message */
.success-banner {
    background: linear-gradient(135deg, #ECFDF5 0%, #D1FAE5 100%);
    border: 1px solid #6EE7B7;
    border-radius: 14px;
    padding: 1.25rem;
    text-align: center;
    margin: 1rem 0;
}

.success-banner h3 {
    color: #065F46;
    margin: 0 0 0.25rem 0;
    font-size: 1.1rem;
}

.success-banner p {
    color: #047857;
    margin: 0;
    font-size: 0.9rem;
}

/* Footer */
.footer-text {
    text-align: center;
    color: #94A3B8;
    font-size: 0.85rem;
    padding: 1.5rem 0;
    margin-top: 1rem;
}
</style>
""", unsafe_allow_html=True)


# ============ MAIN APPLICATION ============

# Hero Section
st.markdown("""
<div class="hero-section">
    <h1 class="hero-title">✨ Brilliant PPT Generator</h1>
    <p class="hero-subtitle">Transform any idea into a CEO-impressive presentation in seconds</p>
    <div class="hero-badge">⚡ Powered by Google Gemini AI - 100% FREE</div>
</div>
""", unsafe_allow_html=True)

# Main Layout
col1, col2 = st.columns([1, 1], gap="large")

with col1:
    # Content Input Card
    st.markdown('<div class="modern-card">', unsafe_allow_html=True)
    st.markdown('''
        <div class="card-header">
            <span class="card-header-icon">📝</span>
            <span class="card-header-text">Your Content</span>
        </div>
    ''', unsafe_allow_html=True)
    
    pres_title = st.text_input(
        "Presentation Title",
        placeholder="e.g., Digital Transformation Strategy 2025",
        help="Enter a clear, professional title for your presentation"
    )
    
    content = st.text_area(
        "Your Ideas / Notes / Bullet Points",
        height=220,
        placeholder="""Enter anything - raw notes, bullet points, or even just keywords!

Examples:
• Market size $50B, growing 15% yearly
• Competitors: Apple, Google, Microsoft
• Our solution: AI-powered, 50% cost reduction
• Expected ROI: 3x in first year

Or just type topics:
sales strategy, customer acquisition, Q4 goals

The AI will transform ANY input into a professional presentation!""",
        help="Enter your content in any format - the AI will structure it professionally"
    )
    
    st.markdown('</div>', unsafe_allow_html=True)
    
    # Features Card
    st.markdown('<div class="modern-card">', unsafe_allow_html=True)
    st.markdown('''
        <div class="card-header">
            <span class="card-header-icon">🤖</span>
            <span class="card-header-text">AI Capabilities</span>
        </div>
    ''', unsafe_allow_html=True)
    
    st.markdown("""
    <div class="feature-list">
        <div class="feature-item">✓ Creates compelling narratives from any input</div>
        <div class="feature-item">✓ Auto-corrects spelling & grammar errors</div>
        <div class="feature-item">✓ Generates professional business language</div>
        <div class="feature-item">✓ Structures content with logical flow</div>
        <div class="feature-item">✓ Works with incomplete or messy notes</div>
        <div class="feature-item">✓ Produces 4-8 optimized slides</div>
    </div>
    """, unsafe_allow_html=True)
    
    # Number of slides
    num_slides = st.slider("Number of Content Slides", min_value=4, max_value=10, value=6)
    
    st.markdown('</div>', unsafe_allow_html=True)

with col2:
    # Theme Selection Card
    st.markdown('<div class="modern-card">', unsafe_allow_html=True)
    st.markdown('''
        <div class="card-header">
            <span class="card-header-icon">🎨</span>
            <span class="card-header-text">Choose Theme</span>
        </div>
    ''', unsafe_allow_html=True)
    
    selected_theme = st.selectbox(
        "Color Theme",
        options=list(THEMES.keys()),
        index=0,
        help="Select a professional color theme"
    )
    
    theme = THEMES[selected_theme]
    
    # Theme preview
    st.markdown(f"""
        <div style="
            background: {theme['gradient']};
            height: 70px;
            border-radius: 14px;
            margin: 0.75rem 0;
            display: flex;
            align-items: center;
            justify-content: center;
            color: white;
            font-weight: 700;
            font-size: 1.1rem;
            box-shadow: 0 4px 15px {theme['preview_dark']}40;
        ">{selected_theme}</div>
    """, unsafe_allow_html=True)
    
    # All themes grid
    st.markdown("**All Available Themes:**")
    theme_cols = st.columns(5)
    for idx, (name, t) in enumerate(THEMES.items()):
        with theme_cols[idx % 5]:
            is_selected = name == selected_theme
            border_style = f"3px solid #1E293B" if is_selected else "1px solid #E2E8F0"
            st.markdown(f"""
                <div class="theme-item">
                    <div class="theme-preview-box" style="background: {t['gradient']}; border: {border_style};"></div>
                    <div class="theme-label">{'✓ ' if is_selected else ''}{name.split()[0]}</div>
                </div>
            """, unsafe_allow_html=True)
    
    st.markdown('</div>', unsafe_allow_html=True)
    
    # Tips Card
    st.markdown('<div class="modern-card">', unsafe_allow_html=True)
    st.markdown('''
        <div class="card-header">
            <span class="card-header-icon">💡</span>
            <span class="card-header-text">Pro Tips</span>
        </div>
    ''', unsafe_allow_html=True)
    
    st.markdown("""
    **Best Results:**
    - Include specific numbers and data
    - Mention your target audience
    - Add key topics you want covered
    - Use bullet points for structure
    
    **AI Will Handle:**
    - Messy or unorganized notes
    - Spelling mistakes
    - Grammar errors
    - Incomplete sentences
    - Missing context
    """)
    st.markdown('</div>', unsafe_allow_html=True)

# Generate Button
st.markdown("---")
col_b1, col_b2, col_b3 = st.columns([1, 2, 1])
with col_b2:
    generate_btn = st.button("🚀 Generate Professional Presentation", type="primary", use_container_width=True)

# Session State
if 'generated_slides' not in st.session_state:
    st.session_state.generated_slides = None
if 'generated_title' not in st.session_state:
    st.session_state.generated_title = None
if 'pptx_data' not in st.session_state:
    st.session_state.pptx_data = None

# Generate Presentation
if generate_btn:
    if not pres_title.strip():
        st.error("⚠️ Please enter a presentation title")
    else:
        with st.spinner(""):
            # Progress indicators
            progress_container = st.container()
            
            with progress_container:
                st.markdown('<div class="progress-step"><span>🔄</span><span class="progress-step-text">Analyzing your content...</span></div>', unsafe_allow_html=True)
                time.sleep(0.5)
                
                st.markdown('<div class="progress-step"><span>🤖</span><span class="progress-step-text">AI is creating your story...</span></div>', unsafe_allow_html=True)
            
            # Generate with AI
            result = generate_presentation_with_ai(pres_title, content, num_slides)
            
            # Retry once if failed
            if not result:
                with progress_container:
                    st.markdown('<div class="progress-step"><span>🔁</span><span class="progress-step-text">Retrying with optimized prompt...</span></div>', unsafe_allow_html=True)
                result = generate_presentation_with_ai(pres_title, content, num_slides)
            
            # Fallback to smart parser
            if not result:
                with progress_container:
                    st.markdown('<div class="progress-step"><span>📝</span><span class="progress-step-text">Using smart content parser...</span></div>', unsafe_allow_html=True)
                result = smart_fallback_parser(pres_title, content)
            
            if result and result.get('slides'):
                st.session_state.generated_slides = result['slides']
                st.session_state.generated_title = result.get('title', pres_title)
                
                with progress_container:
                    st.markdown('<div class="progress-step"><span>📊</span><span class="progress-step-text">Generating PowerPoint file...</span></div>', unsafe_allow_html=True)
                
                # Create PPTX
                st.session_state.pptx_data = create_pptx(
                    st.session_state.generated_title,
                    st.session_state.generated_slides,
                    selected_theme
                )
                
                st.balloons()
                st.success(f"✅ Created {len(st.session_state.generated_slides)} professional slides!")
            else:
                st.error("❌ Generation failed. Please try again.")

# Display Results
if st.session_state.generated_slides:
    st.markdown("---")
    
    # Stats
    st.markdown(f"### 📊 {st.session_state.generated_title}")
    st.markdown(f"""
    <div class="stats-row">
        <div class="stat-badge">📑 {len(st.session_state.generated_slides)} Content Slides</div>
        <div class="stat-badge">🎨 Theme: {selected_theme}</div>
        <div class="stat-badge">📐 16:9 Widescreen</div>
    </div>
    """, unsafe_allow_html=True)
    
    # Slide Previews
    for idx, slide in enumerate(st.session_state.generated_slides):
        with st.expander(f"📄 Slide {idx + 1}: {slide.get('title', 'Untitled')}", expanded=(idx < 3)):
            r, g, b = theme['header_bg']
            bullets_html = ''.join([f'<p>• {bullet}</p>' for bullet in slide.get('bullets', [])])
            st.markdown(f"""
            <div class="slide-card">
                <div class="slide-card-header" style="background: rgb({r},{g},{b});">
                    {slide.get('title', 'Untitled')}
                </div>
                <div class="slide-card-body">
                    {bullets_html}
                </div>
            </div>
            """, unsafe_allow_html=True)
    
    # Download Section
    st.markdown("---")
    
    st.markdown("""
    <div class="success-banner">
        <h3>🎉 Your Presentation is Ready!</h3>
        <p>Click below to download your CEO-impressive PowerPoint</p>
    </div>
    """, unsafe_allow_html=True)
    
    col_d1, col_d2, col_d3 = st.columns([1, 2, 1])
    with col_d2:
        # Create filename
        safe_title = re.sub(r'[^\w\s-]', '', st.session_state.generated_title)
        safe_title = re.sub(r'\s+', '_', safe_title).lower()[:50]
        filename = f"{safe_title}.pptx"
        
        st.download_button(
            label="📥 Download PowerPoint (PPTX)",
            data=st.session_state.pptx_data,
            file_name=filename,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )

# Footer
st.markdown("""
<div class="footer-text">
    Made with ❤️ | Powered by Google Gemini AI | 100% Free Forever
</div>
""", unsafe_allow_html=True)
